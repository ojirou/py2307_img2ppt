import os
from pptx import Presentation
from pptx.util import Inches
# 画像フォルダのパス
image_folder = r'C:\\Users\\user\\python\\ppt\\sample'
# パワーポイントプレゼンテーションを作成
presentation = Presentation()
# スライドに画像を配置する
slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # レイアウト6を使用するスライドを作成
# 画像フォルダ内の画像ファイルを取得
image_files = [f for f in os.listdir(image_folder) if f.endswith(('.jpg', '.jpeg', '.png', '.gif'))]
# スライド上に画像を配置する位置とサイズの設定
left = Inches(1)  # 左端の位置（インチ単位）
top = Inches(1)  # 上端の位置（インチ単位）
width = Inches(2)  # 画像の幅（インチ単位）
space = Inches(0.5)  # 画像間のスペース（インチ単位）
slide_height = presentation.slide_height
max_height = slide_height * 3 / 4  # スライドの高さの3/4までの領域に収める
# 画像をスライドに配置する
for i, image_file in enumerate(image_files):
    image_path = os.path.join(image_folder, image_file)
    image = slide.shapes.add_picture(image_path, left, top, width)
    image_ratio = image.width / image.height
    # 画像の高さを調整してスライドの高さの3/4以内に収める
    if image.height > max_height:
        image.height = max_height
        image.width = max_height * image_ratio
    if (i + 1) % 3 == 0:
        top += image.height + space
        # top += max_height + space
        left = Inches(1)
    else:
        left += image.width + space
# 生成されたパワーポイントファイルの保存パス
output_file = 'output.pptx'
# パワーポイントファイルを保存
presentation.save(output_file)
# 保存したパワーポイントファイルを開く（Windows環境の場合）
os.startfile(output_file)